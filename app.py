import os
import io
import base64
import streamlit as st
from openai import OpenAI

api_key = st.secrets.get("OPENAI_API_KEY") or os.environ.get("OPENAI_API_KEY")
if not api_key:
    st.error("OPENAI_API_KEY is not set. Please add it in Streamlit Secrets.")
    st.stop()
client = OpenAI(api_key=os.environ.get("OPENAI_API_KEY"))

# -------- Agent 1: Job Summary --------
SUMMARY_PROMPT = """
You are an AI assistant specialized in understanding writing tasks and producing a structured Job Summary, not the full content itself. Read the user’s instructions and any extracted text from attachments (e.g., PDFs, DOCX) to identify what needs to be written, including topic, word count or length, reference style (APA, MLA, Harvard, etc.), and writing style or document type (essay, report, PPT, proposal, article, dissertation, thesis, etc.). If a detail is not explicitly given but can be reasonably inferred, infer it; if it cannot be inferred confidently, mark it as “Not specified.” Always respond in this exact format: Topic: <short topic or title>; Word Count: <number of words or If word count is not mentioned in the Job card, then by default print "1500">; Reference Style: <style or If Reference Style is not mentioned in the Job card, then by default print "Harvard">; Writing Style: <type or "Report">; Job Summary: <10–20 sentences clearly describing what needs to be written, the main themes to cover, target audience or level if known, and any important constraints such as tone or structure>. Do not add extra sections, do not explain your reasoning, and do not write the actual assignment—only provide a clear, concise, implementation-ready Job Summary that another writer or AI could directly follow.
"""

# -------- Agent 2: Structure with word breakdown --------
STRUCTURE_PROMPT = """
You are an AI assistant specialized in creating academic writing structures (detailed outlines) for writing tasks. Your input is always the full output of a Job Summary agent, which includes at least: Topic, Word Count, Reference Style, Writing Style, and Job Summary (and may also include extra instructions). Your job is to design a clear, logically ordered, academically appropriate structure with word counts for each section and subsection, so that another writer or AI could directly draft the final document. Strictly follow all instructions and requirements from the Job Summary and ensure that every key theme, focus area, or constraint is reflected in the structure. Use academic writing conventions that match the Writing Style (e.g., essays with introduction/body/conclusion; reports with sections such as introduction, methodology, analysis, conclusion; dissertations/thesis with chapters such as introduction, literature review, methodology, results, discussion, conclusion; PPTs as slide-based academic sections, etc.). Handle Word Count as follows: always use only word counts and never pages, lines, slides, or any other length unit; if a specific word count is given, treat it as the target total and allocate section word counts so they sum to approximately that total (with minor acceptable variation); if a range is given, internally pick a reasonable midpoint and allocate based on that; if the word count is described in pages or similar, internally convert to an approximate word count and output only word counts; if Word Count is “Not specified,” infer a reasonable total based on the Writing Style and academic context, then allocate accordingly. Respect the Reference Style by including a final “References” or “Bibliography” section with an appropriate word count whenever references are expected for that type of task. Ensure a coherent hierarchy with numbered sections and, where useful, subsections, each with a clear academic-style heading and an explicit word count (e.g., “Section Title – X words”). Begin by stating the title (using the Topic) and the total word count, then list the sections in order. Do not write any actual content of the sections, only the structure and word counts. Do not explain your reasoning, do not add extra metadata fields, and do not mention any unit other than words.sub points must shows the word counts in word breakdown and the sum of sub points word count must match with the main points total word count
"""

# -------- Agent 3: Content generation from structure --------
CONTENT_PROMPT = """
You are an AI assistant specialized in academic content writing. Your input is the full output of a Structure-Making Agent, which includes the title, total word count, and a numbered list of sections and subsections with individual word counts, all derived from a Job Summary. Your task is to transform this structure into complete, polished content that strictly follows all instructions, rules, and constraints implied by both the structure and the underlying task (topic, writing style, level, focus areas, tone, etc.). You must: (1) preserve the given headings and their order exactly as provided; (2) write cohesive, formal, academic prose under each section/subsection that clearly addresses the intent of its heading and the overall task; (3) follow the specified word counts closely for each section and subsection, aiming to be as close as reasonably possible to the target for each one and to the overall total; (4) maintain consistency in voice, tense, and perspective as implied by the task; and (5) ensure logical flow between sections with appropriate transitions and internal coherence. Must follow the exact word count that is mentioned, but sometimes you can provide 5% More or less in the contents as word counts. Do not modify or invent new sections, do not change the title, and do not contradict any explicit requirements from the task (such as focus, scope, or audience). When writing the content, do not include any reference list, bibliography, or citations of any kind (no in-text citations, no author-year, no numbers in brackets, and no “References” section), even if the structure or task mentions a reference style; treat that aspect as handled elsewhere. Do not explain your reasoning or describe your process; output only the final written content organized under the given headings.
"""

# -------- Agent 4: References + in-text citations list --------
REFERENCES_PROMPT = """
You are an AI assistant specialized in generating academic reference lists and corresponding in-text citation formats. Your input will be: (1) the full content produced by a content-creation agent, (2) the specified reference style (e.g., APA, MLA, Chicago, Harvard, IEEE, etc.), and (3) the approximate total word count of the content. Your task is to create an original, topic-related reference list that strictly follows the given reference style and is based on the themes, concepts, and topics present in the content. All references you provide must be to real, credible, and verifiable sources published after 2021 (i.e., from 2022 onwards). For every 1000 words of content, generate approximately 7 references (rounding reasonably to the nearest whole number) and ensure that all references are directly relevant to the subject matter of the content. Present the references as a properly formatted “Reference List” ordered alphabetically (A–Z) by the first author’s surname, strictly conforming to the rules of the specified reference style. After the alphabetical reference list, provide a separate “Citation List” that contains the in-text citation format for each reference above (e.g., for Harvard and APA: Author, Year; for MLA: Author page; for IEEE: [number], etc.), covering all references already listed. In-text Citation rules: For Harvard, APA, APA7,  IEEE Referencing (If one, two, or three authors are present in the Reference, then use the Surname of Each Author first, then a comma, and then the year in a Single bracket). Like example: ‘Hermes, A. and Riedl, R., 2021, July. Dimensions of retail customer experience and its outcomes: a literature review and directions for future research. If you notice here, two authors are present, so the in-text citation will be “(Hermes and Riedl, 2021)”. If 4 or more authors are present, then use the first author's surname, then et al., then a comma, and then the year. For example: “Pappas, A., Fumagalli, E., Rouziou, M. and Bolander, W., 2023. More than machines: The role of the future retail salesperson in enhancing the customer experience. Journal of Retailing, 99(4), pp.518-531.”. If you notice here 4 authors are present, so the intext citation will be (Pappas et al. 2023). In IEEE, all are the same but in Number Format like [1], [2], etc. Do not include any explanation, analysis, or extra text beyond the reference list and the citation list. Do not rewrite or summarize the original content. Your entire output must consist only of the formatted reference list followed by the citation list.
"""

# -------- Agent 5: Final document with citations inserted --------
FINALIZE_PROMPT = """
You are an AI assistant specialized in finalizing academic documents by inserting in-text citations and appending an existing reference list. Your inputs are: (1) a complete piece of content with no citations or reference list, (2) a formatted reference list, (3) a citation list that specifies the correct in-text citation format for each reference, and (4) the reference style to follow (e.g., APA, MLA, Chicago, Harvard, IEEE, etc.). Your task is to cite all existing references from the citation list within the content and then append the full reference list at the end of the document, strictly following the given reference style. You must not rewrite, expand, shorten, reorder, or otherwise change any of the existing content, headings, or wording; you may only insert in-text citations at appropriate locations and add the reference list at the end. Don't cite in the Introduction, Conclusion parts, and if available, Abstract and Executive summary; in those parts, don't add in-text citations. Do not add new references, do not remove any existing references, and do not invent sources. Ensure that every reference from the provided reference list is cited at least once in the body using the corresponding in-text format from the citation list, and that all in-text citations match entries in the reference list. Maintain the original structure and formatting of the content as much as possible, only adding the necessary citation markers and the final reference list section. As output, return the full content with the in-text citations properly inserted and the complete reference list appended at the end, and do not include any explanations, notes, or extra commentary.

"""

ALLOWED_EXTENSIONS = [
    "doc", "docx", "pdf",
    "png", "jpg", "jpeg",
    "pptx", "csv", "xlsx", "xlx",
]

# ---------- helper: text extraction ----------

def extract_text_from_upload(uploaded_file) -> str:
    """
    Best-effort text extraction from different file types.
    """
    filename = uploaded_file.name
    name_lower = filename.lower()
    ext = os.path.splitext(name_lower)[1]

    data = uploaded_file.getvalue()

    try:
        if ext == ".pdf":
            from PyPDF2 import PdfReader
            reader = PdfReader(io.BytesIO(data))
            parts = []
            for page in reader.pages:
                text = page.extract_text() or ""
                parts.append(text)
            return "\n".join(parts).strip()

        elif ext == ".docx":
            from docx import Document
            doc = Document(io.BytesIO(data))
            parts = [para.text for para in doc.paragraphs]
            return "\n".join(parts).strip()

        elif ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        parts.append(shape.text)
            return "\n".join(parts).strip()

        elif ext in [".csv"]:
            import pandas as pd
            df = pd.read_csv(io.BytesIO(data))
            return df.to_csv(index=False)

        elif ext in [".xlsx", ".xlx", ".xls"]:
            import pandas as pd
            df = pd.read_excel(io.BytesIO(data))
            return df.to_csv(index=False)

        elif ext == ".doc":
            return (
                "[.doc file detected. Automatic extraction is limited. "
                "Please convert to .docx or PDF for better results.]"
            )

        else:
            try:
                return data.decode("utf-8", errors="ignore")
            except Exception:
                return (
                    f"[Could not automatically extract text from {filename}. "
                    f"Please provide instructions in the text box.]"
                )

    except ImportError as ie:
        return (
            f"[Missing Python library to parse {filename}: {ie}. "
            "Install required libraries (PyPDF2, python-docx, python-pptx, pandas, openpyxl).]"
        )
    except Exception as e:
        return f"[Error while reading {filename}: {e}]"


# ---------- Agent 1: generate job summary ----------

def generate_job_summary(instruction_text, uploaded_files, model="gpt-4.1-mini"):
    attachment_text_blocks = []
    image_contents = []

    for uf in uploaded_files:
        filename = uf.name
        ext = os.path.splitext(filename.lower())[1]

        # Images -> input_image
        if ext in [".png", ".jpg", ".jpeg"]:
            raw = uf.getvalue()
            mime = "image/png" if ext == ".png" else "image/jpeg"
            b64 = base64.b64encode(raw).decode("utf-8")
            image_contents.append(
                {
                    "type": "input_image",
                    "image_url": f"data:{mime};base64,{b64}",
                }
            )
        else:
            text = extract_text_from_upload(uf)
            if text:
                attachment_text_blocks.append(
                    f"----- File: {filename} -----\n{text}"
                )

    base_instruction = (instruction_text or "").strip()
    if not base_instruction:
        base_instruction = (
            "The instructions for the writing task are in the following extracted "
            "file contents. Please infer all possible details about the assignment."
        )

    all_attachments_text = "\n\n".join(attachment_text_blocks).strip()
    if all_attachments_text:
        combined_text = (
            base_instruction
            + "\n\nBelow is the extracted text from the uploaded files:\n\n"
            + all_attachments_text
        )
    else:
        combined_text = base_instruction

    content_items = [{"type": "input_text", "text": combined_text}]
    content_items.extend(image_contents)

    messages = [
        {
            "role": "user",
            "content": content_items,
        }
    ]

    response = client.responses.create(
        model=model,
        instructions=SUMMARY_PROMPT,
        input=messages,
    )

    return response.output_text


# ---------- Agent 2: generate structure from summary ----------

def generate_structure_from_summary(job_summary_text, model="gpt-4.1-mini"):
    messages = [
        {
            "role": "user",
            "content": [
                {"type": "input_text", "text": job_summary_text}
            ],
        }
    ]

    response = client.responses.create(
        model=model,
        instructions=STRUCTURE_PROMPT,
        input=messages,
    )

    return response.output_text


# ---------- Agent 3: generate content from structure ----------

def generate_content_from_structure(structure_text, model="gpt-4.1-mini"):
    messages = [
        {
            "role": "user",
            "content": [
                {"type": "input_text", "text": structure_text}
            ],
        }
    ]

    response = client.responses.create(
        model=model,
        instructions=CONTENT_PROMPT,
        input=messages,
    )

    return response.output_text


# ---------- Agent 4: generate references & in-text citations list ----------

def generate_references_from_content(content_text, reference_style, total_words, model="gpt-4.1"):
    combined = (
        f"Reference style: {reference_style}\n"
        f"Approximate total word count: {total_words}\n\n"
        f"Content:\n{content_text}"
    )

    messages = [
        {
            "role": "user",
            "content": [
                {"type": "input_text", "text": combined}
            ],
        }
    ]

    response = client.responses.create(
        model=model,
        instructions=REFERENCES_PROMPT,
        input=messages,
    )

    return response.output_text


# ---------- Agent 5: finalize document with in-text citations + reference list ----------

def generate_final_document_with_citations(
    content_text,
    reference_list,
    citation_list,
    reference_style,
    model="gpt-4.1"
):
    combined = (
        f"Reference style: {reference_style}\n\n"
        "=== CONTENT (NO CITATIONS) ===\n"
        f"{content_text}\n\n"
        "=== REFERENCE LIST ===\n"
        f"{reference_list}\n\n"
        "=== CITATION LIST ===\n"
        f"{citation_list}\n"
    )

    messages = [
        {
            "role": "user",
            "content": [
                {"type": "input_text", "text": combined}
            ],
        }
    ]

    response = client.responses.create(
        model=model,
        instructions=FINALIZE_PROMPT,
        input=messages,
    )

    return response.output_text


# ---------- Streamlit UI ----------

st.set_page_config(page_title="Click To Assignment", page_icon="📝", layout="centered")

st.title("📝 Click To Assignment (Summary → Structure → Content → References → Final)")

# Session state for all stages
if "job_summary" not in st.session_state:
    st.session_state["job_summary"] = ""
if "structure" not in st.session_state:
    st.session_state["structure"] = ""
if "content" not in st.session_state:
    st.session_state["content"] = ""
if "references" not in st.session_state:
    st.session_state["references"] = ""
if "final_document" not in st.session_state:
    st.session_state["final_document"] = ""

st.write(
    "Step 1: Generate a **Job Summary** from your brief and files.\n\n"
    "Step 2: Generate a **detailed structure with word breakdown** from the Job Summary.\n\n"
    "Step 3: Generate full **academic content** from the structure.\n\n"
    "Step 4: Generate a **Reference List & Citation List** from the content.\n\n"
    "Step 5: Insert **in-text citations** and append the **reference list** to create the final document."
)

# ---------- Step 1: Summary generation ----------
with st.form("job_summary_form"):
    instruction = st.text_area(
        "Additional instructions (optional)",
        placeholder="E.g. 'Please check the attached brief and summarize what needs to be written.'",
        height=120,
    )

    files = st.file_uploader(
        "Upload attachments",
        type=ALLOWED_EXTENSIONS,
        accept_multiple_files=True,
        help="You can upload DOC, DOCX, PDF, PNG, JPG, JPEG, PPTX, CSV, XLSX, XLX files.",
    )

    submitted_summary = st.form_submit_button("① Generate Job Summary")

if submitted_summary:
    # Allow: instructions only, files only, or both.
    has_instruction = bool(instruction and instruction.strip())
    has_files = bool(files)

    if not has_instruction and not has_files:
        st.error("Please either upload at least one file or enter some instructions.")
    else:
        if not os.environ.get("OPENAI_API_KEY"):
            st.error("OPENAI_API_KEY is not set. Please set it in your environment.")
        else:
            try:
                with st.spinner("Generating Job Summary..."):
                    summary_text = generate_job_summary(instruction, files or [])
                st.session_state["job_summary"] = summary_text
                st.session_state["structure"] = ""
                st.session_state["content"] = ""
                st.session_state["references"] = ""
                st.session_state["final_document"] = ""
            except Exception as e:
                st.error(f"Something went wrong during summary generation: {e}")

# ---------- Show Job Summary + Step 2 ----------
if st.session_state["job_summary"]:
    st.subheader("📌 Job Summary (Agent 1 output)")
    edited_summary = st.text_area(
        "You can edit the Job Summary before generating the structure (optional):",
        value=st.session_state["job_summary"],
        height=250,
    )
    st.session_state["job_summary"] = edited_summary

    if st.button("② Generate Detailed Structure (with word breakdown)"):
        if not os.environ.get("OPENAI_API_KEY"):
            st.error("OPENAI_API_KEY is not set. Please set it in your environment.")
        else:
            try:
                with st.spinner("Generating structure from Job Summary..."):
                    structure_text = generate_structure_from_summary(
                        st.session_state["job_summary"]
                    )
                st.session_state["structure"] = structure_text
                st.session_state["content"] = ""
                st.session_state["references"] = ""
                st.session_state["final_document"] = ""
            except Exception as e:
                st.error(f"Something went wrong during structure generation: {e}")

# ---------- Show Structure + Step 3 ----------
if st.session_state["structure"]:
    st.subheader("📚 Detailed Structure (Agent 2 output)")
    edited_structure = st.text_area(
        "You can edit the structure before generating the content (optional):",
        value=st.session_state["structure"],
        height=350,
    )
    st.session_state["structure"] = edited_structure

    if st.button("③ Generate Full Academic Content"):
        if not os.environ.get("OPENAI_API_KEY"):
            st.error("OPENAI_API_KEY is not set. Please set it in your environment.")
        else:
            try:
                with st.spinner("Generating content from structure..."):
                    content_text = generate_content_from_structure(
                        st.session_state["structure"]
                    )
                st.session_state["content"] = content_text
                st.session_state["references"] = ""
                st.session_state["final_document"] = ""
            except Exception as e:
                st.error(f"Something went wrong during content generation: {e}")

# ---------- Show content + Step 4 ----------
if st.session_state["content"]:
    st.subheader("📖 Final Academic Content (Agent 3 output)")
    edited_content = st.text_area(
        "You can edit the content before generating references (optional):",
        value=st.session_state["content"],
        height=400,
    )
    st.session_state["content"] = edited_content

    st.markdown("### 📚 Step 4: Generate Reference List & In-text Citation List")

    default_word_count = len(edited_content.split()) if edited_content.strip() else 0

    col1, col2 = st.columns(2)
    with col1:
        reference_style = st.text_input("Reference style", value="Harvard", key="ref_style_step4")
    with col2:
        total_words = st.number_input(
            "Approximate total word count",
            min_value=0,
            value=default_word_count,
            key="wc_step4",
        )

    if st.button("④ Generate References & Citation List"):
        if not os.environ.get("OPENAI_API_KEY"):
            st.error("OPENAI_API_KEY is not set. Please set it in your environment.")
        else:
            try:
                with st.spinner("Generating references and in-text citations..."):
                    refs_text = generate_references_from_content(
                        st.session_state["content"],
                        reference_style,
                        total_words,
                    )
                st.session_state["references"] = refs_text
                st.session_state["final_document"] = ""
            except Exception as e:
                st.error(f"Something went wrong during reference generation: {e}")

# ---------- Show references + Step 5 ----------
if st.session_state["references"]:
    st.subheader("📚 Reference List & Citation List (Agent 4 output)")
    refs_raw = st.session_state["references"]

    # Try to auto-split into Reference List and Citation List using the heading
    default_ref_list = refs_raw
    default_cit_list = ""
    if "Citation List" in refs_raw:
        head, tail = refs_raw.split("Citation List", 1)
        default_ref_list = head.strip()
        default_cit_list = "Citation List" + tail.strip()

    st.markdown("You can adjust the Reference List and Citation List before finalizing (optional):")

    reference_style_final = st.text_input(
        "Reference style (for final document)",
        value=st.session_state.get("ref_style_step4", "Harvard"),
        key="ref_style_step5",
    )

    reference_list_text = st.text_area(
        "Reference List",
        value=default_ref_list,
        height=250,
        key="ref_list_text",
    )

    citation_list_text = st.text_area(
        "Citation List",
        value=default_cit_list,
        height=250,
        key="cit_list_text",
    )

    if st.button("⑤ Generate Final Document with Citations"):
        if not os.environ.get("OPENAI_API_KEY"):
            st.error("OPENAI_API_KEY is not set. Please set it in your environment.")
        else:
            try:
                with st.spinner("Inserting in-text citations and appending reference list..."):
                    final_doc = generate_final_document_with_citations(
                        content_text=st.session_state["content"],
                        reference_list=reference_list_text,
                        citation_list=citation_list_text,
                        reference_style=reference_style_final,
                    )
                st.session_state["final_document"] = final_doc
            except Exception as e:
                st.error(f"Something went wrong during final document generation: {e}")

# ---------- Show final document ----------
if st.session_state["final_document"]:
    st.subheader("✅ Final Document with In-text Citations and Reference List (Agent 5 output)")
    st.text_area(
        "Final document (copy-paste friendly)",
        value=st.session_state["final_document"],
        height=500,
    )
